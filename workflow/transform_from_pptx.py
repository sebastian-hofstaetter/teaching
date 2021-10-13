#
# Transform a narrated powerpoint file and output:
# - Transcripts per slide (Using Azure speech)
# - Timestamps and slide titles for YouTube's chapter UI
# -------------------------------
#

import argparse
import os
import sys
sys.path.append(os.getcwd())
import copy
import time
import datetime

from pptx import Presentation
from collections import OrderedDict
from pydub import AudioSegment
import azure.cognitiveservices.speech as speechsdk

parser = argparse.ArgumentParser()
parser.add_argument('--in-pptx', action='store', dest='in_file', help='', required=True)
#parser.add_argument('--out-pptx', action='store', dest='out_file', help='bool ', required=True)
parser.add_argument('--out-cc', action='store', dest='out_file_notes', help='', required=True)
parser.add_argument('--azure-key', action='store', dest='azure_key', help='', required=True)
parser.add_argument('--slide-no',nargs='+', action='store',type=int, dest='slide_numbers', help='', required=False)
args = parser.parse_args()

start_time = time.time()

pp = Presentation(args.in_file)

closed_captions = OrderedDict()

run_all_slides = args.slide_numbers is None
slide_number = 0
for slide in pp.slides:

    # skip hidden slides
    if slide._element.get('show') != None and slide._element.get('show') == "0":
        continue

    if not run_all_slides:
        if slide_number not in args.slide_numbers:
            slide_number += 1
            continue

    # hacky way of getting the title 
    title = ""
    all_text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            if shape.name.startswith("Title"):
                title += shape.text + " "
            all_text += shape.text + " "

    words = [word for word in all_text.split() if word.isalpha()]
    # only lowercase a word if its not an acronym with multiple uppercases
    for t,w in enumerate(words):
        sum_uppercase = sum((1 for c in w if c.isupper()))
        if sum_uppercase == 1:
            words[t] = w.lower()
    word_set = list(set(words))

    bi_grams = list(set([word+" "+ words[i+1] for i, word in enumerate(words[:-1])]))
    tri_grams = list(set([word+" "+ words[i+1]+" "+ words[i+2] for i, word in enumerate(words[:-2])]))

    # audio is stored in related_parts
    print("Slide",slide_number+1)
    known_parts = set() 
    for id_, rel in copy.deepcopy(slide.part.related_parts).items():
        if rel.content_type == "audio/mp4" and rel.partname not in known_parts:
            known_parts.add(rel.partname)
            print(rel.partname)

            #
            # extract & convert audio
            #
            with open("temp.m4a","wb") as temp_file:
                temp_file.write(rel.blob)

            # convert and save as azure-speech default audio format 
            track = AudioSegment.from_file("temp.m4a")
            track = track.set_frame_rate(16_000)
            track = track.set_sample_width(2) # 2 byte (16 bit) samples
            track = track.set_channels(1)
            track.export("temp.wav", format='wav')
            track_duration = track.duration_seconds

            #
            # send to speech recognition service
            #

            speech_config = speechsdk.SpeechConfig(subscription=args.azure_key, region="westeurope",speech_recognition_language="en-US")
            audio_config = speechsdk.audio.AudioConfig(filename="temp.wav")
            speech_recognizer = speechsdk.SpeechRecognizer(speech_config=speech_config, audio_config=audio_config)

            # doesn't seem to be working at the moment ... maybe time will improve it 
            # update 2021: yes it does work now :) 
            phrase_list_grammar = speechsdk.PhraseListGrammar.from_recognizer(speech_recognizer)
            for phrase in (word_set+bi_grams+tri_grams):
                phrase_list_grammar.addPhrase(phrase)

            done = False
            recognized_results = ""

            def stop_cb(evt):
                global done
                #print('CLOSING on {}'.format(evt))
                done = True
        
            def recognized_cb(evt):
                global recognized_results
                #print('RECOGNIZED: {}'.format(evt))
                recognized_results += evt.result.text + " "

            # Connect callbacks to the events fired by the speech recognizer
            speech_recognizer.recognized.connect(recognized_cb)
            # stop continuous recognition on either session stopped or canceled events
            speech_recognizer.session_stopped.connect(stop_cb)
            speech_recognizer.canceled.connect(stop_cb)


            # Start continuous speech recognition
            speech_recognizer.start_continuous_recognition()
            while not done:
                time.sleep(.5)

            speech_recognizer.stop_continuous_recognition()

            # cleanup
            phrase_list_grammar.clear()
            del phrase_list_grammar
            del speech_recognizer
            os.remove("temp.m4a")
            os.remove("temp.wav")

            #
            # fill in text as notes
            #
            
            # adding notes breaks powerpoint print to notes ... needs fix
            #slide.notes_slide.notes_text_frame.text = "[Auto-generated CC] " + recognized_results
            closed_captions[slide_number] = (title, recognized_results, track_duration)
            slide_number += 1

with open(args.out_file_notes,"w",encoding="utf8") as markdown_writer:

    markdown_writer.write("# "+os.path.splitext(os.path.basename(args.in_file))[0]+"\n\n")
    markdown_writer.write("*Automatic closed captions generated with the Azure Speech API*\n\n")

    durations=[]
    print("Slide-Setlist: ")
    for slide,(title,cc_text,track_duration) in closed_captions.items():
        markdown_writer.write("### **"+ str(slide+1)+"** "+title+"\n")
        markdown_writer.write(cc_text+"\n\n*" + str(round(track_duration,2)) +" seconds*\n\n")

        print(str(datetime.timedelta(seconds=round(sum(durations),0)))+" "+str(slide+1)+" - " + title)
        durations.append(track_duration)

    if len(durations) > 1:
        markdown_writer.write("### Stats\n")
        markdown_writer.write("Average talking time: " + str(sum(durations)/len(durations)))

print("Finished! After",time.time()-start_time,"seconds")
#pp.save(args.out_file)